#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Batch invitation generator

⚙️ Requires: 
    LibreOffice installed and available in PATH
📄 Usage:
    python generate_invites.py

This script:
1. Reads guest names from a text file.
2. Replaces a placeholder in a PowerPoint template.
3. Exports each customized slide as a JPG.
4. Compresses the final images into a ZIP file.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.exc import PackageNotFoundError
from pathlib import Path
import subprocess


TEMPLATE_FILE = "invite.pptx"
NAMES_FILE = "names.txt"
OUTPUT_DIR = Path("./output-images")
ZIP_NAME = "invitations.zip"

PLACEHOLDER = "{}"

FONT_NAME = "2 Davat"
FONT_SIZE = Pt(36)

def check_dependency(command: str, name: str):
    """Check if a required command is available."""
    try:
        subprocess.run([command, "--version"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except FileNotFoundError:
        print(f"❌ Error: {name} not installed or not in PATH.")
        exit(1)


def load_names(file_path: str) -> list[str]:
    """Return list of guest names from `file_path`."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return [name.strip() for name in f.read().splitlines() if name.strip()]
    except FileNotFoundError:
        print(f"❌ Error: Names file not found → {file_path}")
        return []


def replace_placeholder(slide, name: str):
    """Replace placeholders in slide text with the guest name."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        frame = shape.text_frame
        for p in frame.paragraphs:
            for r in p.runs:
                if PLACEHOLDER in r.text:
                    r.text = r.text.replace(PLACEHOLDER, name)
                    
                    r.font.name = FONT_NAME
                    r.font.size = FONT_SIZE


def save_and_convert(pptx_path: Path, output_path: Path):
    """Convert PPTX file to JPG using LibreOffice and remove the PPTX file."""
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to", "jpg",
        str(pptx_path),
        "--outdir", str(output_path),
    ]
    
    try:
        subprocess.run(cmd, check=True)
        pptx_path.unlink(missing_ok=True)
        
    except FileNotFoundError:
        print("❌ Error: LibreOffice not installed or not in PATH.")
    
    except subprocess.CalledProcessError as e:
        print(f"❌ LibreOffice conversion failed for {pptx_path.name}: {e}")


def create_invites(names: list[str]):
    """Generate invitation slides and export them as images."""
    OUTPUT_DIR.mkdir(exist_ok=True)
    
    total = len(names)
    for i, name in enumerate(names, start=1):
        print(f"🧩 Processing {i}/{total}: {name}")
        
        try:
            prs = Presentation(TEMPLATE_FILE)
            
        except PackageNotFoundError:
            print(f"❌ Error: Template file not found → {TEMPLATE_FILE}")
            return
        
        slide = prs.slides[0]
        replace_placeholder(slide, name)
        
        output_pptx = OUTPUT_DIR / f"دعوت‌نامه-{name.replace(' ', '-')}.pptx"
        prs.save(output_pptx)
        
        save_and_convert(output_pptx, OUTPUT_DIR)


def compress_files(file_path: Path):
    """Compress generated images into a ZIP archive."""
    if not file_path.exists():
        print(f"⚠️ Directory not found: {file_path}")
        return
    
    try:
        subprocess.run(["zip", "-r", ZIP_NAME, str(file_path)], check=True)
    
    except FileNotFoundError:
        print("❌ Error: 'zip' command not found. Please install it.")
        
    except subprocess.CalledProcessError:
        print("❌ Error: Failed to create ZIP archive.")


if __name__ == "__main__":
    try:
        check_dependency("libreoffice", "LibreOffice")
        check_dependency("zip", "ZIP utility")
        
        names = load_names(NAMES_FILE)
        if not names:
            print("⚠️ No names found — exiting.")
            exit(1)
            
        create_invites(names)
        compress_files(OUTPUT_DIR)
        
    except Exception as e:
        print(f"❌ Unexpected error: {e}")
